"""
多格式檔案轉換器 - 支援 Excel、Word、PowerPoint 等格式轉換為 PDF
專為教學型RAG系統設計，自動處理多種Office文件格式
"""

import os
import tempfile
import logging
from pathlib import Path
from typing import Optional, Dict, Any, List
import shutil
import subprocess
import platform

# Office文件處理
try:
    import pandas as pd
    from openpyxl import load_workbook
    from pptx import Presentation
    from docx import Document
    import win32com.client as win32
    WIN32_AVAILABLE = True
except ImportError as e:
    WIN32_AVAILABLE = False
    # 只在 Windows 系統下顯示警告，Linux 下是預期行為
    if platform.system() == "Windows":
        logging.warning(f"Win32 COM 不可用，將使用替代方案: {e}")

# PDF生成
try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    logging.warning("ReportLab 不可用，將使用替代PDF生成方案")

# 設置日誌
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class FileConverter:
    """多格式檔案轉換器"""
    
    def __init__(self):
        self.supported_formats = {
            '.xlsx': self._convert_excel_to_pdf,
            '.xls': self._convert_excel_to_pdf,
            '.docx': self._convert_word_to_pdf,
            '.doc': self._convert_word_to_pdf,
            '.pptx': self._convert_ppt_to_pdf,
            '.ppt': self._convert_ppt_to_pdf,
            '.pdf': self._copy_pdf  # PDF直接複製
        }
        
        # 檢查系統環境
        self.is_windows = platform.system() == "Windows"
        self.use_com = WIN32_AVAILABLE and self.is_windows
        
        logger.info(f"檔案轉換器初始化完成")
        logger.info(f"支援格式: {list(self.supported_formats.keys())}")
        logger.info(f"Windows COM 可用: {self.use_com}")
        logger.info(f"ReportLab 可用: {REPORTLAB_AVAILABLE}")

    def is_supported_format(self, filename: str) -> bool:
        """檢查檔案格式是否支援"""
        ext = Path(filename).suffix.lower()
        return ext in self.supported_formats

    def convert_to_pdf(self, input_path: str, output_path: str = None) -> Optional[str]:
        """
        將支援的檔案格式轉換為PDF
        
        Args:
            input_path: 輸入檔案路徑
            output_path: 輸出PDF路徑（可選）
            
        Returns:
            成功時返回PDF檔案路徑，失敗時返回None
        """
        try:
            if not os.path.exists(input_path):
                logger.error(f"輸入檔案不存在: {input_path}")
                return None
            
            # 獲取檔案副檔名
            ext = Path(input_path).suffix.lower()
            
            if ext not in self.supported_formats:
                logger.error(f"不支援的檔案格式: {ext}")
                return None
            
            # 生成輸出路徑
            if output_path is None:
                output_path = str(Path(input_path).with_suffix('.pdf'))
            
            logger.info(f"開始轉換檔案: {input_path} -> {output_path}")
            
            # 執行轉換
            converter_func = self.supported_formats[ext]
            success = converter_func(input_path, output_path)
            
            if success and os.path.exists(output_path):
                logger.info(f"檔案轉換成功: {output_path}")
                return output_path
            else:
                logger.error(f"檔案轉換失敗: {input_path}")
                return None
                
        except Exception as e:
            logger.error(f"檔案轉換過程中發生錯誤: {e}")
            return None

    def _copy_pdf(self, input_path: str, output_path: str) -> bool:
        """直接複製PDF檔案"""
        try:
            if input_path != output_path:
                shutil.copy2(input_path, output_path)
            return True
        except Exception as e:
            logger.error(f"PDF複製失敗: {e}")
            return False

    def _convert_excel_to_pdf(self, input_path: str, output_path: str) -> bool:
        """轉換Excel檔案為PDF"""
        try:
            # 優先使用COM（Windows）
            if self.use_com:
                return self._convert_excel_with_com(input_path, output_path)
            else:
                return self._convert_excel_with_pandas(input_path, output_path)
        except Exception as e:
            logger.error(f"Excel轉換失敗: {e}")
            return False

    def _convert_excel_with_com(self, input_path: str, output_path: str) -> bool:
        """使用COM轉換Excel（Windows專用）"""
        try:
            excel = win32.Dispatch("Excel.Application")
            excel.Visible = False
            excel.DisplayAlerts = False
            
            workbook = excel.Workbooks.Open(os.path.abspath(input_path))
            workbook.ExportAsFixedFormat(0, os.path.abspath(output_path))
            
            workbook.Close()
            excel.Quit()
            return True
        except Exception as e:
            logger.error(f"COM Excel轉換失敗: {e}")
            return False

    def _convert_excel_with_pandas(self, input_path: str, output_path: str) -> bool:
        """使用pandas和reportlab轉換Excel"""
        if not REPORTLAB_AVAILABLE:
            logger.error("ReportLab不可用，無法轉換Excel")
            return False
            
        try:
            # 讀取Excel檔案
            if input_path.endswith('.xlsx'):
                df_dict = pd.read_excel(input_path, sheet_name=None, engine='openpyxl')
            elif input_path.endswith('.xls'):
                # 對於舊版 .xls 檔案，使用 xlrd 引擎
                df_dict = pd.read_excel(input_path, sheet_name=None, engine='xlrd')
            else:
                # 自動偵測
                df_dict = pd.read_excel(input_path, sheet_name=None)
            
            # 創建PDF
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            story = []
            styles = getSampleStyleSheet()
            
            for sheet_name, df in df_dict.items():
                # 工作表標題
                story.append(Paragraph(f"工作表: {sheet_name}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                # 轉換DataFrame為表格
                if not df.empty:
                    data = [df.columns.tolist()] + df.values.tolist()
                    table = Table(data)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black)
                    ]))
                    story.append(table)
                    story.append(Spacer(1, 12))
            
            doc.build(story)
            return True
            
        except Exception as e:
            logger.error(f"Pandas Excel轉換失敗: {e}")
            logger.error(f"檔案路徑: {input_path}")
            logger.error(f"檔案副檔名: {Path(input_path).suffix}")

            # 提供具體的解決建議
            if "xlrd" in str(e).lower():
                logger.error("建議: 請安裝 xlrd 套件來處理 .xls 檔案: pip install xlrd==2.0.1")
            elif "openpyxl" in str(e).lower():
                logger.error("建議: 請安裝 openpyxl 套件來處理 .xlsx 檔案: pip install openpyxl")
            elif "reportlab" in str(e).lower():
                logger.error("建議: 請安裝 reportlab 套件來生成 PDF: pip install reportlab")

            return False

    def _convert_word_to_pdf(self, input_path: str, output_path: str) -> bool:
        """轉換Word檔案為PDF"""
        try:
            # 優先使用COM（Windows）
            if self.use_com:
                return self._convert_word_with_com(input_path, output_path)
            else:
                return self._convert_word_with_docx(input_path, output_path)
        except Exception as e:
            logger.error(f"Word轉換失敗: {e}")
            return False

    def _convert_word_with_com(self, input_path: str, output_path: str) -> bool:
        """使用COM轉換Word（Windows專用）"""
        try:
            word = win32.Dispatch("Word.Application")
            word.Visible = False
            
            doc = word.Documents.Open(os.path.abspath(input_path))
            doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # 17 = PDF格式
            
            doc.Close()
            word.Quit()
            return True
        except Exception as e:
            logger.error(f"COM Word轉換失敗: {e}")
            return False

    def _convert_word_with_docx(self, input_path: str, output_path: str) -> bool:
        """使用python-docx和reportlab轉換Word"""
        if not REPORTLAB_AVAILABLE:
            logger.error("ReportLab不可用，無法轉換Word")
            return False
            
        try:
            doc = Document(input_path)
            
            # 創建PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            story = []
            styles = getSampleStyleSheet()
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    story.append(Paragraph(paragraph.text, styles['Normal']))
                    story.append(Spacer(1, 6))
            
            pdf_doc.build(story)
            return True
            
        except Exception as e:
            logger.error(f"python-docx Word轉換失敗: {e}")
            return False

    def _convert_ppt_to_pdf(self, input_path: str, output_path: str) -> bool:
        """轉換PowerPoint檔案為PDF"""
        try:
            # 優先使用COM（Windows）
            if self.use_com:
                return self._convert_ppt_with_com(input_path, output_path)
            else:
                # 在Linux環境下，先嘗試LibreOffice，再嘗試python-pptx
                if self._convert_ppt_with_libreoffice(input_path, output_path):
                    return True
                else:
                    return self._convert_ppt_with_pptx(input_path, output_path)
        except Exception as e:
            logger.error(f"PowerPoint轉換失敗: {e}")
            return False

    def _convert_ppt_with_libreoffice(self, input_path: str, output_path: str) -> bool:
        """使用LibreOffice轉換PowerPoint（Linux/Mac）"""
        try:
            import subprocess

            # 檢查LibreOffice是否可用
            try:
                subprocess.run(['libreoffice', '--version'],
                             capture_output=True, check=True, timeout=10)
            except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
                logger.warning("LibreOffice不可用，跳過此轉換方法")
                return False

            # 創建輸出目錄
            output_dir = os.path.dirname(output_path)
            os.makedirs(output_dir, exist_ok=True)

            # 使用LibreOffice轉換
            cmd = [
                'libreoffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', output_dir,
                input_path
            ]

            logger.info(f"執行LibreOffice轉換: {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)

            if result.returncode == 0:
                # LibreOffice會生成與輸入檔案同名的PDF
                input_name = os.path.splitext(os.path.basename(input_path))[0]
                generated_pdf = os.path.join(output_dir, f"{input_name}.pdf")

                # 如果生成的PDF名稱與期望的不同，重新命名
                if generated_pdf != output_path and os.path.exists(generated_pdf):
                    shutil.move(generated_pdf, output_path)

                if os.path.exists(output_path):
                    logger.info(f"LibreOffice轉換成功: {output_path}")
                    return True
                else:
                    logger.error("LibreOffice轉換後找不到輸出檔案")
                    return False
            else:
                logger.error(f"LibreOffice轉換失敗: {result.stderr}")
                return False

        except Exception as e:
            logger.error(f"LibreOffice PowerPoint轉換失敗: {e}")
            return False

    def _convert_ppt_with_com(self, input_path: str, output_path: str) -> bool:
        """使用COM轉換PowerPoint（Windows專用）"""
        try:
            ppt = win32.Dispatch("PowerPoint.Application")
            ppt.Visible = 1

            presentation = ppt.Presentations.Open(os.path.abspath(input_path))
            presentation.SaveAs(os.path.abspath(output_path), 32)  # 32 = PDF格式

            presentation.Close()
            ppt.Quit()
            return True
        except Exception as e:
            logger.error(f"COM PowerPoint轉換失敗: {e}")
            return False

    def _convert_ppt_with_pptx(self, input_path: str, output_path: str) -> bool:
        """使用python-pptx和reportlab轉換PowerPoint"""
        if not REPORTLAB_AVAILABLE:
            logger.error("ReportLab不可用，無法轉換PowerPoint")
            return False
            
        try:
            prs = Presentation(input_path)
            
            # 創建PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            story = []
            styles = getSampleStyleSheet()
            
            for i, slide in enumerate(prs.slides):
                # 投影片標題
                story.append(Paragraph(f"投影片 {i+1}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                # 提取文字內容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        story.append(Paragraph(shape.text, styles['Normal']))
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 20))
            
            pdf_doc.build(story)
            return True
            
        except Exception as e:
            logger.error(f"python-pptx PowerPoint轉換失敗: {e}")
            return False

    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """獲取檔案資訊"""
        try:
            path = Path(file_path)
            return {
                "filename": path.name,
                "extension": path.suffix.lower(),
                "size": os.path.getsize(file_path),
                "supported": self.is_supported_format(path.name),
                "exists": path.exists()
            }
        except Exception as e:
            logger.error(f"獲取檔案資訊失敗: {e}")
            return {}

if __name__ == "__main__":
    # 測試範例
    converter = FileConverter()
    
    # 測試檔案轉換
    test_files = [
        "test.xlsx",
        "test.docx", 
        "test.pptx"
    ]
    
    for test_file in test_files:
        if os.path.exists(test_file):
            print(f"\n測試轉換: {test_file}")
            result = converter.convert_to_pdf(test_file)
            if result:
                print(f"轉換成功: {result}")
            else:
                print("轉換失敗")
        else:
            print(f"測試檔案不存在: {test_file}")
